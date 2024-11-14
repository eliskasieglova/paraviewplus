# powerpoint maker
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from graphmaker import TimeSeriesDemonstration



class PresentationMaker:
    def __init__(self):

        self.output_folder = ""
        self.name = ""

        self.ppt = Presentation()

        self.main_logo = "paraviewplus/ppt_utils/main_logo.png"
        self.logo = "paraviewplus/ppt_utils/logo.png"
        
        self.sw = 10
        self.ppt.slide_width = Inches(self.sw)
        self.ppt.slide_height = Inches(5.62519685)

        # bounds for text
        self.left = 0.787402
        self.text_width = 8.4094488

        self.logo_width = 0.472441

    def set_output_folder(self, output_folder):
        self.output_folder = output_folder

    def set_name(self, name):
        self.name = name

    def save(self):
        self.ppt.save(self.output_folder + "/" + self.name + ".pptx")

    def change_font(self, font_name):
        for slide in self.ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):  # Ensure the shape has text
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = font_name  # Set font family

    def add_banner(self, slide):

        # Green-to-blue gradient banner on top
        shape = slide.shapes.add_shape(1, 0, 0, self.ppt.slide_width, Inches(0.5))  # Rectangle shape
        fill = shape.fill
        fill.gradient()  # Set the fill type to gradient
        fill.gradient_angle = 45  # make the gradient horizontal

        # Define the gradient stops (green to blue)
        stop1 = fill.gradient_stops[0]
        stop1.position = 0.0
        stop1.color.rgb = RGBColor(149, 193, 39)  # Green color
        stop2 = fill.gradient_stops[1]
        stop2.position = 1.0
        stop2.color.rgb = RGBColor(70, 172, 226)  # Blue color

        # Remove border and shadow
        shape.line.fill.background()  # Removes the border by setting it to background
        shape.shadow.inherit = False  # Disable shadow

    def add_bottom_logo(self, slide):
        # place logo
        logo = slide.shapes.add_picture(self.logo, Inches(self.sw / 2 - self.logo_width / 2), Inches(5), width=Inches(self.logo_width))

    def add_universal_title(self, slide):
        # style title
        title = slide.shapes.title

        title.left, title.top, title.height, title.width = Inches(self.left), Inches(0.9409449), Inches(0.5866142), Inches(self.text_width)
        title_para = title.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(23)
        title_para.alignment = PP_ALIGN.LEFT
        
        return title

    def create_title_slide(self):
        # slide layout
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[0])
        
        # add gradient banner on top
        self.add_banner(slide)

        # place logo
        logo = slide.shapes.add_picture(self.main_logo, Inches(0.787402), Inches(0.629921), Inches(1.433071))

        # style title
        title = slide.shapes.title
        title.left, title.top, title.height, title.width = Inches(0.787402), Inches(2), Inches(1.677165), Inches(8.4094488)
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(38)
        title_para.font.color.rgb = RGBColor(51, 51, 51)
        title_para.alignment = PP_ALIGN.LEFT

        # style subtitle
        subtitle = slide.placeholders[1]
        subtitle.left, subtitle.top, subtitle.height, subtitle.width = Inches(0.787402), Inches(3.720472), Inches(0.590551), Inches(8.4094488)
        subtitle_para = subtitle.text_frame.paragraphs[0]
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = RGBColor(51, 51, 51)
        subtitle_para.alignment = PP_ALIGN.LEFT
        

        # TODO background picture, ecoten logo

        return slide
    
    def create_content_slide(self):
        # normal text slide
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[1])

        # add banner to top
        self.add_banner(slide)

        # add title
        title = self.add_universal_title(slide)

        # style paragraph
        paragraph = slide.shapes.placeholders[1]

        self.paragraph_height = 3.38583

        paragraph.left, paragraph.top, paragraph.height, paragraph.width = Inches(self.left), Inches(1.8), Inches(self.paragraph_height), Inches(self.text_width)  # Change the width here
        para_frame = paragraph.text_frame
        para_frame.paragraphs[0].font.size = Pt(13)
        para_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
        para_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        self.add_bottom_logo(slide)

        return slide

    def create_comparison_slide(self):
        # split text slide
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[3])

        # add banner to top
        self.add_banner(slide)

        # add logo
        self.add_bottom_logo(slide)

        # add title
        title = self.add_universal_title(slide)

        # style paragraphs
        left_paragraph = slide.shapes.placeholders[1]

        left_paragraph.left = Inches(self.left)
        left_paragraph.top = Inches(1.8)
        left_paragraph.height = Inches(self.paragraph_height)
        left_paragraph.width = Inches(self.text_width / 2 - self.text_width * 0.006)  # Change the width here
        left_para_frame = left_paragraph.text_frame.paragraphs[0]
        left_para_frame.font.size = Pt(13)
        left_para_frame.font.color.rgb = RGBColor(51, 51, 51)
        left_para_frame.alignment = PP_ALIGN.LEFT

                # style paragraphs
        right_paragraph = slide.shapes.placeholders[2]

        right_paragraph.left = Inches(self.left + self.text_width / 2 + self.text_width * 0.006)
        right_paragraph.top = Inches(1.8)
        right_paragraph.height = Inches(self.paragraph_height)
        right_paragraph.width = Inches(self.text_width / 2 - self.text_width * 0.006)  # Change the width here
        right_para_frame = right_paragraph.text_frame.paragraphs[0]
        right_para_frame.font.size = Pt(13)
        right_para_frame.font.color.rgb = RGBColor(51, 51, 51)
        right_para_frame.alignment = PP_ALIGN.LEFT

        return slide

    def create_chapter_slide(self, c):
        # title slide for chapters
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[5])

        # background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(70, 172, 226) if c == "blue" else RGBColor(149, 193, 39) 

        # add title
        title = slide.shapes.title
        title.left, title.top, title.height, title.width = Inches(self.left), Inches(1.1), Inches(1.677165), Inches(self.text_width)
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # add line
        line = slide.shapes.add_shape(1, left=Inches(self.left + self.text_width * 0.05 / 2), top=Inches(1), width=Inches(self.text_width - self.text_width * 0.05), height=Inches(0.06)
        )

        line.line.color.rgb = RGBColor(255, 255, 255)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
        line.line.fill.background()  # Removes the border by setting it to background
        line.shadow.inherit = False  # Disable shadow

        # add logo
        logo = slide.shapes.add_picture(f"paraviewplus/ppt_utils/logo_{c}.png", Inches(self.sw / 2 - self.logo_width / 2), Inches(5), width=Inches(self.logo_width))

        return slide

    def create_title_only_slide(self):
        # title only
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[5])

        # add top banner
        self.add_banner(slide)

        #add bottom logo
        self.add_bottom_logo(slide)

        title = self.add_universal_title(slide)

        return slide
    
    def create_two_content_slide(self):

        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[7])

        # background color
        left_half = slide.shapes.add_shape(
            1,  # Shape type: 1 corresponds to a rectangle
            0, 0,  # Position the shape at the top-left corner
            self.ppt.slide_width / 2, self.ppt.slide_height  # Width is half the slide, full height
        )
        left_half.fill.solid()  # Fill with solid color
        left_half.fill.fore_color.rgb = RGBColor(238, 245, 219)  # Green color
        left_half.line.fill.background()  # Removes the border by setting it to background
        left_half.shadow.inherit = False  # Disable shadow

        slide.shapes._spTree.remove(left_half._element)
        slide.shapes._spTree.insert(2, left_half._element)  # move to background

        # layout - banner, logo
        self.add_banner(slide)
        self.add_bottom_logo(slide)

        # adjust title 
        title = slide.shapes.title
        title.left, title.top, title.height, title.width = Inches(0.787402), Inches(1.440945), Inches(1.846457), Inches(3.610236)
        title_para = title.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(26)
        title_para.alignment = PP_ALIGN.LEFT

        # adjust box 1 
        box1 =  slide.shapes.placeholders[1]
        box1.left, box1.top, box1.height, box1.width = Inches(0.787402), Inches(3.456693), Inches(0.8307087), Inches(3.610236)
        box1_para = box1.text_frame.paragraphs[0]
        box1_para.font.size = Pt(16)
        box1_para.alignment = PP_ALIGN.LEFT

        # adjust box 2
        box2 = slide.shapes.placeholders[2]
        box2.left, box2.top, box2.height, box2.width = Inches(5.4094488), Inches(1.480315), Inches(3.30709), Inches(3.688976)
        box2_para = box1.text_frame.paragraphs[0]
        box2_para.font.size = Pt(16)
        box2_para.alignment = PP_ALIGN.LEFT


        return slide

    def create_blank_slide(self):

        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[6])

        self.add_banner(slide)

        self.add_bottom_logo(slide)

        return slide
    
        
    def create_footnote_slide(self):

        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[6])  # blank layout
        # layout
        self.add_banner(slide)
        self.add_bottom_logo(slide)
        
        # add text
        footnote = slide.shapes.add_textbox(Inches(self.left), Inches(3.976378), Inches(self.text_width), Inches(0.7))
        footnote_para = footnote.text_frame.add_paragraph()
        footnote_para.text = "Kliknutím vložíte text."
        footnote_para.font.size = Pt(13)
        footnote_para.font.color.rgb = RGBColor(0, 0, 0)

        return slide

    def create_end_slide(self):

        # title slide for chapters
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[1])

        # background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(70, 172, 226)

        # add title
        title = slide.shapes.title
        title.left, title.top, title.height, title.width = Inches(self.left), Inches(self.left), Inches(1.362205), Inches(self.text_width)
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.LEFT

        # style subtitle
        subtitle = slide.placeholders[1]
        subtitle.left, subtitle.top, subtitle.height, subtitle.width = Inches(self.left), Inches(2.484252), Inches(1.728346), Inches(self.text_width)
        subtitle_para = subtitle.text_frame.paragraphs[0]
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = RGBColor(51, 51, 51)
        subtitle_para.alignment = PP_ALIGN.LEFT

        # add line
        line = slide.shapes.add_shape(1, left=Inches(0.8937008), top=Inches(4.5511811), width=Inches(0.8070866), height=Inches(0.06)
        )

        line.line.color.rgb = RGBColor(255, 255, 255)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
        line.line.fill.background()  # Removes the border by setting it to background
        line.shadow.inherit = False  # Disable shadow


        # add logo
        logo = slide.shapes.add_picture(f"paraviewplus/ppt_utils/logo_blue.png", Inches(self.sw / 2 - self.logo_width / 2), Inches(5), width=Inches(self.logo_width))


        return slide

    
    def create_template(self):

        self.create_title_slide()
        self.create_content_slide()
        self.create_comparison_slide()
        self.create_chapter_slide("blue")
        self.create_title_only_slide()
        self.create_two_content_slide()
        self.create_chapter_slide("green")

        self.create_blank_slide()
        self.create_footnote_slide()
        self.create_end_slide()

        self.change_font("Proxima Nova")

        self.save()

    def time_series_demonstration(self):

        slides = []
        
        slide1 = self.create_chapter_slide("blue")
        title = slide1.shapes.title
        title.text = "Time Series Demonstration"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        import os
        time_series_folder = "paraviewplus/figs/timeseries/"
        figs = os.listdir(time_series_folder)
        fig = figs[0].split("_")[0]
        
        for idx in range(len(figs)):
            fig_path = time_series_folder + fig + "_" + str(idx + 1) + ".png"
            slide = self.create_blank_slide()
            slide.shapes.add_picture(fig_path, top=Inches(0.5), left=Inches(2.30315), height=Inches(5.3976378))
            
    def create_presentation(self, time_series_demonstration=True):

        self.create_title_slide()

        if time_series_demonstration:

            self.time_series_demonstration()

        self.change_font("Proxima Nova")

        self.save()


#POWERPOINT PRESENTATION
from pptmaker import PresentationMaker

ppt = PresentationMaker()
ppt.set_output_folder("paraviewplus")
ppt.set_name("timeseries_test")
ppt.create_presentation(time_series_demonstration=True)
